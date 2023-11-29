from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation


def read_pptx(path_to_presentation: str) -> list[str]:
    """
            Extracts the text from each shape in each slide of a PowerPoint presentation.

            Args:
                path_to_presentation (str): The file path to the PowerPoint presentation.

            Returns:
               List of string of the slides.
            """
    prs = Presentation(path_to_presentation)

    slide_text = []
    for slide in prs.slides:
        text = ""
        for shape in slide.shapes:
            text += extract_text_from_shape(shape)
        if text.strip():
            slide_text.append(text.strip())

    return slide_text


def extract_text_from_shape(shape) -> str:
    """
            Extracts the text from a shape, including its nested shapes.

            Args:
                shape: The shape object.

            Returns:
                str: The extracted text.
            """
    text = ""
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text += "".join(run.text.strip() for run in paragraph.runs)
            text += "\n"
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            text += extract_text_from_shape(sub_shape)
    return text