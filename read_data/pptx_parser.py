from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
from typing import List


def read_pptx(path_to_presentation: str) -> List[str]:
    """
    Reads a PowerPoint presentation file and extracts text from each shape in each slide.
    Args:
        path_to_presentation (str): The file path to the PowerPoint presentation.
    Returns:
        List[str]: A list of strings, each containing the extracted text from a slide.
    """
    prs = Presentation(path_to_presentation)
    slides_text = []
    for slide in prs.slides:
        text = ""
        for shape in slide.shapes:
            text += extract_text_from_shape(shape)
        if text.strip():
            slides_text.append(text.strip())
    return slides_text


def extract_text_from_shape(shape) -> str:
    """
    Extracts text content from a PowerPoint shape, including nested shapes, tables, and groups.
    Args:
        shape: The shape object.
    Returns:
        str: The extracted text content.
    """
    text = ""
    # Extract text directly from the shape if available
    if hasattr(shape, "text") and shape.text.strip():
        text += shape.text.strip() + "\n"
    # Extract text from the text frame of the shape
    elif shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text += " ".join(run.text.strip() for run in paragraph.runs) + "\n"
    # Extract text from a table within the shape
    elif shape.has_table:
        for row in shape.table.rows:
            text += ",".join(cell.text_frame.text.strip() for cell in row.cells) + "\n"
    # Extract text from nested shapes within a group
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            sub_text = extract_text_from_shape(sub_shape)
            if sub_text.strip():
                text += sub_text
    return text
