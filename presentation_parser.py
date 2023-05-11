from pptx import Presentation


class PresentationParser:
    @staticmethod
    def extract_text(path_to_presentation):
        """
        Extracts the text from each shape in each slide of a PowerPoint presentation.

        Args:
            path_to_presentation (str): The file path to the PowerPoint presentation.

        Returns:
           List of string of the slides.
        """
        prs = Presentation(path_to_presentation)

        slide_text = []
        for slide in prs.slides:
            text = ""
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text += "".join(run.text.strip() for run in paragraph.runs)
                    text += "\n"
            if text.strip():
                slide_text.append(text.strip())

        return slide_text
